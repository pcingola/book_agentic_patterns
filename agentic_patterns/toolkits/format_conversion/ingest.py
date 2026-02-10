"""Ingest: convert PDF/Word/PPT/Excel to Markdown or CSV."""

from pathlib import Path

import pymupdf
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation


def docx_to_markdown(input_path: Path) -> str:
    """Convert a Word document to Markdown."""
    doc = DocxDocument(input_path)
    lines: list[str] = []
    for para in doc.paragraphs:
        style = para.style.name or ""
        text = para.text
        if not text.strip():
            lines.append("")
            continue
        if style.startswith("Heading 1"):
            lines.append(f"# {text}")
        elif style.startswith("Heading 2"):
            lines.append(f"## {text}")
        elif style.startswith("Heading 3"):
            lines.append(f"### {text}")
        elif style.startswith("Heading"):
            level = min(int(style.split()[-1]) if style.split()[-1].isdigit() else 4, 6)
            lines.append(f"{'#' * level} {text}")
        elif style.startswith("List"):
            lines.append(f"- {text}")
        else:
            lines.append(text)
    return "\n".join(lines)


def pdf_to_markdown(input_path: Path) -> str:
    """Convert a PDF to Markdown using pymupdf."""
    doc = pymupdf.open(str(input_path))
    parts: list[str] = []
    for page in doc:
        parts.append(page.get_text("text"))
    doc.close()
    return "\n\n".join(parts)


def pptx_to_markdown(input_path: Path) -> str:
    """Convert a PowerPoint presentation to Markdown."""
    prs = Presentation(input_path)
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        lines.append("")
    return "\n".join(lines)


def xlsx_to_csv(input_path: Path, sheet: str | None = None) -> str:
    """Convert an Excel sheet to CSV text."""
    wb = load_workbook(input_path, read_only=True, data_only=True)
    ws = wb[sheet] if sheet else wb.active
    rows: list[str] = []
    for row in ws.iter_rows(values_only=True):
        cells = [str(c) if c is not None else "" for c in row]
        rows.append(",".join(cells))
    wb.close()
    return "\n".join(rows)


def xlsx_to_markdown(input_path: Path, sheet: str | None = None) -> str:
    """Convert an Excel sheet to a Markdown table."""
    wb = load_workbook(input_path, read_only=True, data_only=True)
    ws = wb[sheet] if sheet else wb.active
    all_rows: list[list[str]] = []
    for row in ws.iter_rows(values_only=True):
        all_rows.append([str(c) if c is not None else "" for c in row])
    wb.close()
    if not all_rows:
        return ""
    header = all_rows[0]
    lines = [" | ".join(header), " | ".join("---" for _ in header)]
    for row in all_rows[1:]:
        lines.append(" | ".join(row))
    return "\n".join(lines)
